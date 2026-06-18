from pathlib import Path

import pandas as pd
from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "pisi_2026"
FIG = ROOT / "artifacts" / "model" / "v4" / "reports" / "figures"


PROJECT_SHORT = "Danang UrbanAgent AI"
PROJECT_VI = "Danang UrbanAgent AI: Nền tảng agent đô thị kết nối trải nghiệm du lịch và vị trí kinh doanh tại Đà Nẵng"
PROJECT_EN = "Danang UrbanAgent AI: An Urban Agent Platform for Travel Experience Planning and Business Location Intelligence"


def read_stats():
    gg = pd.read_csv(ROOT / "data" / "poi_data_ggmap.csv")
    foody = pd.read_csv(ROOT / "data" / "poi_data_foody.csv")
    metrics = pd.read_csv(ROOT / "artifacts" / "model" / "v4" / "reports" / "metrics" / "training_loss_v4.csv")
    best = metrics[metrics["is_best_model"].astype(str).str.lower().eq("yes")].iloc[0]
    return {
        "gg_rows": len(gg),
        "gg_categories": gg["Category"].nunique(),
        "gg_districts": gg["District"].nunique(),
        "foody_rows": len(foody),
        "foody_categories": foody["Category"].nunique(),
        "foody_districts": foody["District"].nunique(),
        "best_epoch": int(best["epoch"]),
        "recall5": float(best["test_recall_5"]),
        "test_loss": float(best["test_loss"]),
        "test_silhouette": float(best["test_silhouette"]),
    }


def set_doc_style(doc):
    style = doc.styles["Normal"]
    style.font.name = "Times New Roman"
    style.font.size = Pt(14)
    for section in doc.sections:
        section.top_margin = Inches(0.8)
        section.bottom_margin = Inches(0.8)
        section.left_margin = Inches(0.9)
        section.right_margin = Inches(0.9)


def add_heading(doc, text, level=1):
    p = doc.add_heading(text, level=level)
    for run in p.runs:
        run.font.name = "Times New Roman"
    return p


def add_para(doc, text="", bold_prefix=None):
    p = doc.add_paragraph()
    if bold_prefix:
        r = p.add_run(bold_prefix)
        r.bold = True
        r.font.name = "Times New Roman"
        r.font.size = Pt(14)
        r2 = p.add_run(text)
        r2.font.name = "Times New Roman"
        r2.font.size = Pt(14)
    else:
        r = p.add_run(text)
        r.font.name = "Times New Roman"
        r.font.size = Pt(14)
    p.paragraph_format.space_after = Pt(6)
    return p


def add_bullets(doc, items):
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        r = p.add_run(item)
        r.font.name = "Times New Roman"
        r.font.size = Pt(14)


def add_business_model_canvas(doc):
    headers = [
        "Đối tác chính",
        "Hoạt động chính",
        "Giá trị cung cấp",
        "Quan hệ KH",
        "Khách hàng",
    ]
    row1 = [
        "Trường/nhóm nghiên cứu; điểm du lịch/quán địa phương; đơn vị dữ liệu bản đồ; vườn ươm khởi nghiệp; đối tác vận chuyển/đặt dịch vụ.",
        "Thu thập dữ liệu POI; huấn luyện embedding đa phương thức; xây agent điều phối; phát triển hệ chuyên gia chỉ đường, lập lịch trình và phân tích vị trí kinh doanh.",
        "Một urban agent cho Đà Nẵng: khách tìm trải nghiệm phù hợp, người kinh doanh tìm khu vực phù hợp; cả hai dùng chung dữ liệu nhu cầu, POI, tuyến đường và ngữ cảnh địa phương.",
        "Demo tự phục vụ; hội thoại với agent; chỉnh sửa lịch trình; phản hồi trong app; báo cáo giải thích cho người kinh doanh.",
        "Du khách/người dân đi chơi tại Đà Nẵng; nhóm bạn/gia đình; chủ quán/startup địa phương; đối tác du lịch, ăn uống và vận chuyển.",
    ]
    headers2 = [
        "Nguồn lực chính",
        "Kênh phân phối",
        "Cơ cấu chi phí",
        "Dòng doanh thu",
        "Tác động xã hội",
    ]
    row2 = [
        "Dữ liệu Google Maps/Foody; mô hình CLIP + ResNet; backend Node/Python; frontend demo; luật chuyên gia giao thông/lịch trình; bộ nhớ hồ sơ người dùng.",
        "Mobile/web app; API; workshop du lịch thông minh; cộng đồng chủ quán địa phương; mạng xã hội và truyền thông trường.",
        "Máy chủ GPU/CPU; lưu trữ dữ liệu; nhân sự AI/full-stack; vận hành cloud; bản đồ/định tuyến; tích hợp công cụ đặt xe/đặt bàn.",
        "Gói app freemium; itinerary nâng cao; dashboard mật độ khách/nhu cầu cho người bán; API cho đối tác; quảng bá địa điểm minh bạch.",
        "Giúp du khách đi đúng nơi; giúp người bán chọn đúng vị trí và hiểu nhu cầu; phân phối lợi ích du lịch tới các điểm địa phương ở Đà Nẵng.",
    ]
    table = doc.add_table(rows=4, cols=5)
    table.style = "Table Grid"
    for col, text in enumerate(headers):
        table.cell(0, col).text = text
        table.cell(1, col).text = row1[col]
    for col, text in enumerate(headers2):
        table.cell(2, col).text = text
        table.cell(3, col).text = row2[col]


def create_docx(stats):
    doc = Document()
    set_doc_style(doc)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    r = title.add_run("PHỤ LỤC 1\nBẢN THUYẾT MINH DỰ ÁN PISI 2026")
    r.bold = True
    r.font.name = "Times New Roman"
    r.font.size = Pt(16)

    add_para(doc, f"Tên dự án: {PROJECT_VI}")
    add_para(doc, f"Project title: {PROJECT_EN}")
    add_para(doc, "Đội thi: [Tên đội thi]")
    add_para(doc, "Trưởng nhóm: [Họ và tên] - Lớp: [Lớp] - SĐT: [Số điện thoại] - Email: [Email]")
    add_para(doc, "Thành viên: [Thành viên 1], [Thành viên 2], [Thành viên 3] (không quá 05 người)")
    add_para(doc, "Người hướng dẫn: [Tên giảng viên/đơn vị/ vai trò]")

    doc.add_page_break()
    add_heading(doc, "Tóm tắt dự án", 1)
    add_para(
        doc,
        "Dự án xây dựng nền tảng agent đô thị cho Đà Nẵng, kết nối hai nhu cầu dùng chung một lõi AI: khách/người dân tìm địa điểm, lộ trình, lịch trình phù hợp; người kinh doanh tìm khu vực mở dịch vụ dựa trên mật độ khách, nhu cầu và cạnh tranh.",
        "Ý tưởng chính: ",
    )
    add_para(
        doc,
        "Sản phẩm kết hợp dữ liệu vị trí, danh mục, đánh giá, hình ảnh, review và hệ luật địa phương để tạo một agent có thể hiểu ý định, lập kế hoạch, gọi công cụ, ghi nhớ sở thích và giải thích quyết định. "
        "Khách có thể yêu cầu 'tối nay đi cafe yên tĩnh gần biển rồi ăn hải sản', chỉnh lịch trình bằng cách thêm/bớt quán, yêu cầu đặt xe/đặt bàn; người bán có thể hỏi 'mở cafe học bài ở đâu tại Đà Nẵng' và nhận phân tích khu vực, mật độ khách, cạnh tranh, tuyến đường và tệp khách phù hợp.",
        "Mô tả sản phẩm/dịch vụ: ",
    )

    add_heading(doc, "I. Tổng quan đề án - Business Model Canvas", 1)
    add_business_model_canvas(doc)

    add_heading(doc, "II. Mô tả thêm về sản phẩm, dịch vụ", 1)
    add_heading(doc, "1. Tính cần thiết của sản phẩm dịch vụ", 2)
    add_bullets(
        doc,
        [
            "Dự án đã có nguyên mẫu gồm backend Node.js/Python, mô hình multimodal encoder, dữ liệu POI, báo cáo huấn luyện, API gợi ý địa điểm và các module hệ chuyên gia phục vụ chỉ đường.",
            "Du khách thường phải tự ghép nhiều nguồn: bản đồ, review, mạng xã hội, thời tiết, khoảng cách và đặt dịch vụ; người kinh doanh lại thiếu dữ liệu nhu cầu khách, mật độ cạnh tranh và đặc điểm từng khu vực.",
            "Khách hàng quan trọng nhất gồm hai role liên kết: người đi chơi/du khách cần lịch trình phù hợp, và chủ quán/startup địa phương cần chọn vị trí kinh doanh dựa trên dữ liệu.",
            "Giá trị xã hội: giúp du khách khám phá Đà Nẵng sâu hơn, hỗ trợ quán/điểm đến địa phương tiếp cận đúng khách, và biến dữ liệu đô thị thành hạ tầng ra quyết định cho thành phố.",
        ],
    )

    add_heading(doc, "2. Tính khả thi", 2)
    add_bullets(
        doc,
        [
            f"Dữ liệu hiện có: {stats['gg_rows']} POI Google Maps với {stats['gg_categories']} danh mục và {stats['foody_rows']} POI Foody với {stats['foody_categories']} danh mục.",
            "Hệ thống đã có luồng suy luận: người dùng nhập mô tả/ảnh, encoder sinh vector, backend trả Top-K địa điểm tương đồng; hệ chuyên gia có thể mở rộng để xếp tuyến đường, lịch trình và chấm điểm vị trí kinh doanh.",
            f"Mô hình V4 đạt Recall@5 trên tập kiểm thử Foody khoảng {stats['recall5'] * 100:.2f}% tại epoch tốt nhất, cho thấy khả năng thích nghi miền từ Google Maps sang Foody.",
            "Chi phí triển khai giai đoạn đầu thấp vì có thể chạy CPU/GPU nhỏ, tiền xử lý embedding theo lô và mở rộng theo API khi có người dùng.",
        ],
    )

    add_heading(doc, "3. Tính độc đáo, sáng tạo", 2)
    add_bullets(
        doc,
        [
            "Điểm mới không nằm ở việc 'có agent', mà ở lớp abstraction cho tác vụ đô thị Đà Nẵng: hiểu intent, lập plan, gọi tool, nhớ ngữ cảnh, giải thích và chuyển tín hiệu nhu cầu thành insight kinh doanh.",
            "Cùng một lõi embedding phục vụ hai chiều: khách tìm nơi hợp mình; người bán tìm khu vực có nhóm khách hợp concept kinh doanh.",
            "Mô hình V4 sử dụng ResNet cho geometry, CLIP cho text/image và gated fusion; phía agent dùng nhiều công cụ như POI retrieval, route expert, itinerary planner, demand heatmap, booking/action connector.",
            "Định hướng nâng cấp dùng semantic-aware hard negative mining để tách rõ các danh mục dễ nhiễu như cafe, nhà hàng, ăn vặt, điểm tham quan, mua sắm và giải trí đêm.",
        ],
    )

    add_heading(doc, "4. Kế hoạch sản xuất, kinh doanh", 2)
    add_bullets(
        doc,
        [
            "Giai đoạn 1: hoàn thiện dữ liệu, chuẩn hóa danh mục, tăng chất lượng retrieval theo ngữ nghĩa và đóng gói demo web với hai role Khách và Người kinh doanh.",
            "Giai đoạn 2: xây Travel Agent cho phép tạo/sửa lịch trình, thêm quán, đổi tuyến, gợi ý đặt xe/đặt bàn; xây Business Agent phân tích mật độ khách, cạnh tranh và tiềm năng khu vực.",
            "Giai đoạn 3: cung cấp app đô thị, API itinerary/business intelligence và dashboard nhu cầu khách cho đối tác du lịch/địa điểm.",
            "Rủi ro chính gồm dữ liệu nhiễu, thiếu ảnh, lệch miền giữa nền tảng và thông tin thời gian thực; giải pháp là tiền xử lý danh mục, hard negative mining, hậu kiểm semantic penalty, human feedback và tích hợp nguồn dữ liệu cập nhật.",
        ],
    )

    add_heading(doc, "5. Kết quả tiềm năng của dự án", 2)
    add_bullets(
        doc,
        [
            "Nguồn thu: app freemium, gói lập lịch trình nâng cao, dashboard mật độ khách/nhu cầu cho người kinh doanh, API truy vấn POI/lịch trình, quảng bá địa điểm có kiểm soát.",
            "Tăng trưởng: mở rộng từ Đà Nẵng sang Huế, Hội An, TP.HCM; mở rộng từ dữ liệu ăn uống ban đầu sang điểm tham quan, sự kiện, lưu trú, mua sắm và dịch vụ.",
            "Tác động: giúp du khách tiết kiệm thời gian lập kế hoạch, giúp người bán giảm rủi ro chọn vị trí sai, phân phối lượt ghé đến các điểm địa phương và tạo lợi ích trực tiếp cho Đà Nẵng.",
        ],
    )

    add_heading(doc, "6. Nguồn lực thực hiện", 2)
    add_bullets(
        doc,
        [
            "Nhân lực: nhóm phát triển AI/data/backend/frontend, có khả năng huấn luyện PyTorch và triển khai API.",
            "Tài sản sẵn có: mã nguồn mô hình, dữ liệu POI, biểu đồ đánh giá, hệ luật chuyên gia và prototype web.",
            "Nguồn lực cần huy động: cloud/GPU thử nghiệm, chi phí khảo sát người dùng, thiết kế giao diện và truyền thông.",
            "Đối tác mong muốn: giảng viên hướng dẫn, vườn ươm, cộng đồng du lịch địa phương, điểm đến/quán nhỏ, đơn vị cung cấp dữ liệu bản đồ.",
        ],
    )

    add_heading(doc, "7. Các kênh truyền thông", 2)
    add_bullets(
        doc,
        [
            "Kênh trường/CLB khởi nghiệp, Facebook/TikTok demo lịch trình ngắn, workshop cho sinh viên và cộng đồng du lịch địa phương.",
            "Landing page/web demo để người dùng nhập ý tưởng và nhận kết quả ngay.",
            "Case study theo danh mục: cà phê, ăn vặt, nhà hàng gia đình để chứng minh hiệu quả bằng ví dụ thực tế.",
        ],
    )

    add_heading(doc, "Minh chứng kỹ thuật", 1)
    add_para(doc, f"Best epoch V4: {stats['best_epoch']}; Test Recall@5: {stats['recall5'] * 100:.2f}%; Test loss: {stats['test_loss']:.4f}.")
    for name in ["loss_curve_v4.png", "recall_at_5_curve_v4.png", "tsne_clusters_v4.jpg"]:
        path = FIG / name
        if path.exists():
            doc.add_picture(str(path), width=Inches(5.8))

    out = OUT / "PISI_2026_ban_thuyet_minh_du_an.docx"
    doc.save(out)
    return out


def add_slide(prs, title, bullets=None, image=None, footer=f"{PROJECT_SHORT} | PISI 2026"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(246, 248, 251)

    box = slide.shapes.add_textbox(PptInches(0.55), PptInches(0.35), PptInches(12.2), PptInches(0.75))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = PptPt(30)
    p.font.name = "Arial"
    p.font.color.rgb = RGBColor(23, 42, 69)

    left_w = 6.3 if image else 11.8
    body = slide.shapes.add_textbox(PptInches(0.7), PptInches(1.35), PptInches(left_w), PptInches(5.35))
    tf = body.text_frame
    tf.word_wrap = True
    if bullets:
        tf.clear()
        for i, item in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = PptPt(20)
            p.font.name = "Arial"
            p.font.color.rgb = RGBColor(33, 37, 41)
            p.space_after = PptPt(8)

    if image and Path(image).exists():
        slide.shapes.add_picture(str(image), PptInches(7.25), PptInches(1.4), width=PptInches(5.45))

    foot = slide.shapes.add_textbox(PptInches(0.55), PptInches(7.05), PptInches(12), PptInches(0.3))
    p = foot.text_frame.paragraphs[0]
    p.text = footer
    p.font.size = PptPt(10)
    p.font.color.rgb = RGBColor(90, 99, 110)
    return slide


def create_pptx(stats):
    prs = Presentation()
    prs.slide_width = PptInches(13.333)
    prs.slide_height = PptInches(7.5)

    add_slide(
        prs,
        PROJECT_SHORT,
        [
            "Nền tảng agent đô thị cho Đà Nẵng: khách tìm trải nghiệm phù hợp, người kinh doanh tìm vị trí phù hợp.",
            "Lõi chung: POI embedding đa phương thức, hệ chuyên gia định tuyến, itinerary planner và business location intelligence.",
            "Đội thi: [Tên đội thi] | PISI 2026",
        ],
    )
    add_slide(
        prs,
        "Vấn đề",
        [
            "Khách đi chơi phải tự ghép bản đồ, review, mạng xã hội, thời tiết, tuyến đường và đặt dịch vụ.",
            "Người kinh doanh thiếu dữ liệu về mật độ khách, nhu cầu theo khu vực, cạnh tranh và tệp khách phù hợp.",
            "Hai bài toán này liên kết với nhau nhưng hiện thường bị xử lý rời rạc.",
        ],
    )
    add_slide(
        prs,
        "Giải pháp",
        [
            "Travel role: agent gợi ý quán/địa điểm, tạo lịch trình, cho phép thêm quán, đổi tuyến, gợi ý đặt xe/đặt bàn.",
            "Business role: agent gợi ý vị trí kinh doanh, phân tích mật độ khách, cạnh tranh, concept phù hợp và tuyến tiếp cận.",
            "Cùng một lõi dữ liệu tạo vòng lặp: nhu cầu khách -> tín hiệu thị trường -> gợi ý kinh doanh tốt hơn.",
        ],
    )
    add_slide(
        prs,
        "Điểm Mới",
        [
            "Agent không mới; cái mới là abstraction tác vụ đô thị cho Đà Nẵng.",
            "Chuẩn thao tác chung: Intent -> Plan -> Tools -> Route -> Memory -> Market Signal.",
            "Thay vì mỗi chức năng tự xử lý riêng, agent điều phối các tool POI, routing, lịch trình, booking và phân tích kinh doanh.",
        ],
    )
    add_slide(
        prs,
        "Dữ liệu",
        [
            f"Google Maps: {stats['gg_rows']} POI, {stats['gg_categories']} danh mục.",
            f"Foody: {stats['foody_rows']} POI, {stats['foody_categories']} danh mục.",
            "Mỗi POI gồm tên, danh mục, vị trí, đánh giá, giá, ảnh và văn bản tổng hợp cho LLM/CLIP.",
        ],
    )
    add_slide(
        prs,
        "Công nghệ lõi",
        [
            "Version 4: ResNet trích xuất geometry, CLIP trích xuất text và image.",
            "Gated Fusion học trọng số từng modality để phù hợp với truy vấn văn bản hoặc ảnh.",
            "Agent layer điều phối POI retrieval, route expert, itinerary planner, booking/action connector và business insight.",
        ],
    )
    add_slide(
        prs,
        "Kết quả mô hình V4",
        [
            f"Best epoch: {stats['best_epoch']}.",
            f"Test Recall@5 trên Foody: {stats['recall5'] * 100:.2f}%.",
            f"Test loss: {stats['test_loss']:.4f}.",
            "Kết quả cho thấy mô hình có khả năng chuyển miền từ Google Maps sang Foody.",
        ],
        FIG / "recall_at_5_curve_v4.png",
    )
    add_slide(
        prs,
        "Minh họa embedding",
        [
            "t-SNE/UMAP giúp quan sát cấu trúc cụm POI trong không gian biểu diễn.",
            "Hướng cải tiến: semantic-aware hard negative để tách quán cà phê, quán nhậu, nhà hàng khi chúng gần nhau về vị trí.",
        ],
        FIG / "tsne_clusters_v4.jpg",
    )
    add_slide(
        prs,
        "Sản phẩm demo",
        [
            "Role Khách: nhập nhu cầu/ảnh, nhận lịch trình, chỉnh sửa bằng cách thêm quán, đổi thứ tự, chọn phương tiện.",
            "Role Người kinh doanh: nhập concept, nhận khu vực phù hợp, mật độ khách, đối thủ, lý do và rủi ro.",
            "Expert System: phân tích mật độ POI, thời tiết, tuyến đường, cấm rẽ, đường một chiều và quy tắc di chuyển.",
        ],
    )
    add_slide(
        prs,
        "Khách hàng và thị trường",
        [
            "Khách hàng B2C: du khách trẻ, nhóm bạn/gia đình, sinh viên, người dân muốn khám phá Đà Nẵng.",
            "Khách hàng B2B: chủ quán, startup địa phương, điểm đến nhỏ, đơn vị vận chuyển/du lịch.",
            "Lợi thế: tập trung Đà Nẵng, hiểu tuyến đường và nhu cầu địa phương sâu hơn app bản đồ chung.",
        ],
    )
    add_slide(
        prs,
        "Mô hình kinh doanh",
        [
            "Gói miễn phí: demo giới hạn để thu người dùng và phản hồi.",
            "B2C: itinerary nâng cao, lịch trình nhóm, gợi ý đặt xe/đặt bàn qua đối tác.",
            "B2B: dashboard mật độ khách, phân tích vị trí, API POI/lịch trình và quảng bá địa điểm minh bạch.",
        ],
    )
    add_slide(
        prs,
        "Lộ trình",
        [
            "0-3 tháng: chuẩn hóa dữ liệu Đà Nẵng, role switch Khách/Người kinh doanh, retrieval nhanh.",
            "3-6 tháng: Travel Agent có chỉnh sửa lịch trình; Business Agent có heatmap nhu cầu và cạnh tranh.",
            "6-12 tháng: tích hợp đặt xe/đặt bàn, dữ liệu thời gian thực và dashboard đối tác.",
        ],
    )
    add_slide(
        prs,
        "Thông điệp",
        [
            f"{PROJECT_SHORT} biến dữ liệu đô thị Đà Nẵng thành một agent biết lập kế hoạch và hành động.",
            "Một lõi chung phục vụ cả người đi chơi và người kinh doanh.",
            "Giá trị nằm ở sản phẩm xây trên agent: trải nghiệm tốt hơn cho khách, quyết định tốt hơn cho người bán.",
        ],
    )

    out = OUT / "PISI_2026_slide_thuyet_minh.pptx"
    prs.save(out)
    return out


def create_markdown_files(stats):
    form_answers = f"""# Câu trả lời form đăng ký PISI 2026

## Tên dự án (Tiếng Việt + Tiếng Anh)
{PROJECT_VI} / {PROJECT_EN}

## Ý tưởng dự án (tóm tắt 1-2 dòng)
Xây dựng nền tảng agent đô thị cho Đà Nẵng với hai role liên kết: khách/người dân được gợi ý quán, lộ trình và lịch trình cá nhân hóa; người kinh doanh được gợi ý vị trí mở dịch vụ dựa trên mật độ khách, nhu cầu và cạnh tranh. Hệ thống kết hợp học biểu diễn POI đa phương thức, hệ chuyên gia chỉ đường và agent có khả năng lập kế hoạch, gọi công cụ, ghi nhớ ngữ cảnh.

## Mô tả dự án (không quá 200 từ)
{PROJECT_SHORT} là nền tảng agent đô thị tập trung cho Đà Nẵng, kết nối hai bài toán vốn liên quan nhưng thường bị tách rời: khách muốn tìm nơi đi chơi phù hợp và người kinh doanh muốn chọn vị trí mở dịch vụ. Dự án chuẩn hóa dữ liệu POI từ Google Maps/Foody gồm danh mục, vị trí, đánh giá, hình ảnh và bình luận. Mô hình Version 4 dùng ResNet, CLIP và Gated Fusion để sinh embedding đa phương thức cho địa điểm. Trên lõi đó, agent hiểu yêu cầu tự nhiên, lập kế hoạch, gọi công cụ gợi ý POI, chỉ đường, chỉnh sửa lịch trình, gợi ý đặt xe/đặt bàn và ghi nhớ sở thích. Với role người kinh doanh, hệ thống phân tích mật độ khách, cạnh tranh, cụm POI, tuyến tiếp cận và concept phù hợp. Điểm mới là vòng lặp dữ liệu hai chiều: nhu cầu của khách tạo tín hiệu thị trường, còn vị trí kinh doanh tốt hơn làm trải nghiệm đô thị tốt hơn.

## Link Google Drive lưu hồ sơ
[Dán link thư mục Google Drive sau khi upload 3 file: bản thuyết minh, video clip, slide thuyết minh. Nhớ bật quyền “Anyone with the link can view”.]
"""

    video_script = f"""# Kịch bản video thuyết minh dự án PISI 2026

Thời lượng đề xuất: 4 phút 30 giây. Định dạng: quay người thuyết trình kết hợp demo màn hình, sơ đồ agent, biểu đồ mô hình và phụ đề ngắn.

## 0:00 - 0:20 | Mở đầu
Hình ảnh: bản đồ Đà Nẵng, các điểm tham quan/quán địa phương, giao diện demo hai role.
Lời thoại: "Ở Đà Nẵng, khách du lịch luôn hỏi: hôm nay đi đâu, ăn gì, di chuyển thế nào cho hợp gu? Còn người kinh doanh lại hỏi: nên mở quán ở đâu, khu nào có khách phù hợp, cạnh tranh ra sao? Nhóm chúng em xây dựng {PROJECT_VI}, một nền tảng agent đô thị kết nối hai bài toán này bằng cùng một lõi dữ liệu và AI."

## 0:20 - 0:50 | Giới thiệu đội thi
Hình ảnh: thành viên đứng cạnh màn hình demo.
Lời thoại: "Chúng em là đội [Tên đội], gồm [Tên thành viên]. Dự án xuất phát từ một quan sát đơn giản: trải nghiệm của khách và quyết định của người bán không tách rời nhau. Nhu cầu khách tạo nên thị trường, còn vị trí kinh doanh tốt tạo nên trải nghiệm đô thị tốt hơn."

## 0:50 - 1:30 | Vấn đề và ý nghĩa xã hội
Hình ảnh: so sánh tự ghép lịch trình thủ công và nhận gợi ý bằng AI.
Lời thoại: "Hiện nay khách thường phải mở nhiều ứng dụng cùng lúc: bản đồ, review, mạng xã hội, thời tiết, đặt xe. Trong khi đó người kinh doanh lại phải tự đoán mật độ khách và chọn vị trí bằng cảm tính. {PROJECT_SHORT} giải quyết bằng một agent có thể hiểu intent, lập plan, gọi tool, ghi nhớ ngữ cảnh và giải thích quyết định."

## 1:30 - 2:15 | Sản phẩm và công nghệ
Hình ảnh: sơ đồ Intent -> Plan -> Tools -> Route -> Memory -> Market Signal.
Lời thoại: "Điểm mới của dự án không phải chỉ là có agent. Theo góc nhìn của chúng em, giá trị nằm ở lớp abstraction phía trên agent: biến các tác vụ đô thị lặp lại thành một chuẩn chung. Người dùng nói nhu cầu, agent lập kế hoạch, gọi các công cụ POI retrieval, chỉ đường, lịch trình, đặt xe hoặc phân tích kinh doanh, sau đó ghi nhớ phản hồi để phục vụ tốt hơn."

## 2:15 - 3:05 | Demo hai role
Hình ảnh: quay màn hình role Khách và role Người kinh doanh.
Lời thoại: "Ở role Khách, người dùng có thể nhập 'tối nay muốn cafe yên tĩnh gần biển rồi ăn hải sản', agent tạo lịch trình, vẽ tuyến, cho phép thêm quán, đổi thứ tự và gợi ý đặt xe. Ở role Người kinh doanh, người dùng nhập 'mở cafe học bài ở Đà Nẵng', agent gợi ý khu vực, phân tích mật độ khách, đối thủ, cụm trường học, tuyến đường và lý do phù hợp."

## 3:05 - 3:40 | Kết quả và lợi thế cạnh tranh
Hình ảnh: biểu đồ Recall@5, t-SNE, loss curve.
Lời thoại: "Trên tập kiểm thử Foody, mô hình V4 đạt Recall@5 khoảng {stats['recall5'] * 100:.2f}% tại epoch tốt nhất. Lợi thế của dự án là không chỉ lọc từ khóa, mà học một không gian biểu diễn chung giữa văn bản, ảnh, vị trí và danh mục. Cùng không gian đó phục vụ cả hai chiều: khách tìm nơi hợp mình, người bán tìm nơi có khách hợp concept."

## 3:40 - 4:15 | Tính khả thi và mô hình triển khai
Hình ảnh: gói sản phẩm, dashboard, báo cáo.
Lời thoại: "Dự án đã có prototype, dữ liệu {stats['gg_rows']} POI Google Maps và {stats['foody_rows']} POI Foody. Mô hình triển khai gồm app freemium cho khách, gói lịch trình nâng cao, dashboard mật độ khách và phân tích vị trí cho người kinh doanh, API cho đối tác địa phương. Thị trường đầu tiên là Đà Nẵng vì nhóm có dữ liệu và hiểu ngữ cảnh địa phương."

## 4:15 - 4:30 | Kết thúc
Hình ảnh: logo/tên dự án và thông điệp.
Lời thoại: "{PROJECT_SHORT} hướng tới một agent đô thị biết lập kế hoạch và hành động cho Đà Nẵng: khách đi đúng nơi, người bán chọn đúng vị trí, và thành phố khai thác dữ liệu địa phương thông minh hơn. Chúng em xin cảm ơn."
"""

    form_path = OUT / "PISI_2026_cau_tra_loi_form.md"
    video_path = OUT / "PISI_2026_kich_ban_video.md"
    form_path.write_text(form_answers, encoding="utf-8")
    video_path.write_text(video_script, encoding="utf-8")
    return form_path, video_path


def create_roadmap_file(stats):
    roadmap = f"""# Roadmap phát triển {PROJECT_SHORT}

## Định vị sản phẩm mới
{PROJECT_SHORT} không chỉ là hệ thống tìm POI tương đồng, cũng không chỉ là app du lịch. Đây là nền tảng agent đô thị cho Đà Nẵng với hai role dùng chung một lõi: khách/người dân tìm trải nghiệm phù hợp; người kinh doanh tìm vị trí phù hợp. Lõi kỹ thuật hiện có gồm dữ liệu POI, multimodal encoder, backend inference và hệ chuyên gia. Hướng phát triển tiếp theo là xây một lớp abstraction chuẩn cho tác vụ đô thị: `Intent -> Plan -> Tools -> Route -> Memory -> Market Signal`.

Điểm khác biệt theo tinh thần "giá trị nằm ở sản phẩm xây phía trên agent": agent không chỉ trả lời, mà biết chia việc, gọi công cụ, nhớ sở thích, chỉnh sửa kế hoạch, giải thích quyết định và biến tương tác của khách thành insight nhu cầu cho người bán.

## 1. Hướng phát triển Frontend
- Role switch rõ ràng: "Tôi đi chơi" và "Tôi kinh doanh".
- Role Khách: nhập sở thích tự nhiên như "đi chơi tối nay gần biển, ngân sách vừa phải", tải ảnh vibe, nhận lịch trình và bản đồ.
- Chỉnh sửa bằng hội thoại: thêm quán, đổi thứ tự, bỏ điểm quá xa, đổi phương tiện, yêu cầu đặt xe/đặt bàn.
- Role Người kinh doanh: nhập concept như "mở cafe học bài", xem heatmap nhu cầu, mật độ khách, đối thủ, cụm POI và khu vực đề xuất.
- Chế độ giải thích: mỗi gợi ý cần có lý do ngắn như "phù hợp vì gần biển, rating cao, cùng cụm ăn tối và ít lệch tuyến" hoặc "khu vực này có nhiều khách sinh viên nhưng ít cafe học bài".

## 2. Hướng phát triển Backend
- Tách API thành các service rõ ràng: agent orchestrator, POI retrieval, route expert system, itinerary planner, business location scorer, user/business profile.
- Tiền xử lý embedding POI theo batch để inference nhanh, không encode toàn bộ dữ liệu mỗi lần truy vấn.
- Thêm cache theo truy vấn phổ biến và cache tuyến đường theo cặp tọa độ.
- Chuẩn hóa schema POI chung cho Google Maps, Foody và nguồn mới: id, name, category, lat/lon, opening_hours, price, rating, images, reviews, tags.
- API agent nên trả về cả `plan_steps`, `tool_calls`, `score`, `reason`, `warnings`, `route_segments`, `estimated_time`, `fallback_options`, `business_insights`.

## 3. Hướng phát triển mô hình AI
- Nâng cấp Version 4 thành semantic-first retrieval: nếu truy vấn có danh mục rõ ràng, text/category embedding được ưu tiên hơn geometry.
- Huấn luyện semantic-aware hard negative mining: các POI gần nhau nhưng khác intent/danh mục phải bị đẩy xa trong embedding space.
- Thêm reranker nhẹ sau cosine similarity: kết hợp semantic score, category match, distance, rating, opening status, diversity, demand signal và competition penalty.
- Xây dựng bộ đánh giá riêng cho hai role: Recall@K theo intent, itinerary satisfaction, route efficiency, category purity, demand-fit score và business-location score.
- Mở rộng dữ liệu ngoài F&B: điểm tham quan, bãi biển, bảo tàng, sự kiện, lưu trú, shopping, khu vui chơi, điểm check-in.

## 4. Agent và hệ chuyên gia
- Agent Orchestrator: nhận intent, chia việc thành bước nhỏ, chọn tool, tổng hợp kết quả và hỏi lại khi thiếu thông tin.
- Travel Planner Agent: tạo lịch trình, chỉnh sửa lịch trình, thêm quán, đổi tuyến, gợi ý đặt xe/đặt bàn.
- Route Expert Agent: kiểm tra đường một chiều, cấm rẽ, đường cấm theo giờ, tuyến quá vòng, thời gian di chuyển.
- Business Insight Agent: phân tích khu vực, mật độ khách, đối thủ, tệp khách, concept-fit và rủi ro mở dịch vụ.
- Memory/Profile: ghi nhớ sở thích khách, lịch sử chỉnh sửa, concept kinh doanh và phản hồi để lần sau gợi ý tốt hơn.
- Luật định tuyến: tránh đường một chiều sai chiều, cấm rẽ, đường cấm theo giờ, tuyến quá vòng, tuyến qua khu hay kẹt.
- Luật lịch trình: không xếp điểm đóng cửa, không nhảy quá xa giữa hai điểm, ưu tiên cụm địa lý, cân bằng ăn uống/tham quan/nghỉ.
- Luật cá nhân hóa: gia đình có trẻ em ưu tiên ít di chuyển; nhóm bạn ưu tiên check-in/ăn uống; người lớn tuổi ưu tiên điểm nhẹ nhàng.
- Kết hợp scoring: `final_score = semantic_score + route_score + context_score + diversity_score - penalty`.
- Cần giải thích được quyết định, ví dụ: "điểm này được chọn vì cùng hướng với tuyến biển và mở cửa trong khung giờ bạn chọn".

## 5. Tính năng khác biệt có thể đưa vào demo PISI
- "Tạo lịch trình 1 chạm": nhập một câu, app sinh lịch trình nửa ngày/một ngày.
- "Chỉnh lịch trình bằng agent": thêm quán, đổi tuyến, chọn phương tiện, yêu cầu đặt xe/đặt bàn.
- "Business Site Agent": nhập concept kinh doanh, nhận khu vực phù hợp, heatmap khách, đối thủ và lý do.
- "Đi theo vibe": tải ảnh quán/địa điểm, app tìm các nơi có phong cách tương tự.
- "Route-aware recommendation": không chỉ gợi ý điểm hay, mà gợi ý điểm không làm tuyến đường bị vòng.
- "Demand loop": nhu cầu tìm kiếm của khách tạo tín hiệu thị trường ẩn danh cho role người kinh doanh.
- "Explainable Urban Agent": giải thích từng điểm đến hoặc vị trí kinh doanh bằng ngôn ngữ tự nhiên.

## 6. Lộ trình đề xuất
- 0-1 tháng: chuẩn hóa tên dự án, UI hai role, schema POI và API retrieval nhanh bằng precomputed embedding.
- 1-3 tháng: hoàn thiện agent orchestrator, route expert system, bản đồ, route scoring và giải thích tuyến.
- 3-6 tháng: xây Travel Planner Agent, Business Insight Agent, semantic reranker, user/business profile và thử nghiệm tại Đà Nẵng.
- 6-12 tháng: tích hợp đặt xe/đặt bàn, dữ liệu thời gian thực, thêm nguồn POI mới, phát hành beta app và làm việc với đối tác địa phương.

## 7. Các chỉ số nên đo
- Retrieval: Recall@5/10, category purity, semantic mismatch rate.
- Lịch trình: tổng thời gian di chuyển, số điểm phù hợp, tỷ lệ điểm đang mở cửa, độ đa dạng danh mục.
- Người dùng: tỷ lệ chọn lại gợi ý, thời gian tạo lịch trình, điểm hài lòng sau chuyến đi.
- Kinh doanh: demand-fit score, competition score, mật độ khách theo khu vực, số báo cáo vị trí tạo ra, số đối tác địa điểm/dữ liệu.

## 8. Câu chuyện nên dùng khi thuyết minh
Nền tảng ban đầu chứng minh khả năng hiểu POI bằng AI đa phương thức trên dữ liệu Google Maps/Foody: {stats['gg_rows']} POI Google Maps, {stats['foody_rows']} POI Foody, V4 đạt Recall@5 khoảng {stats['recall5'] * 100:.2f}%. Bước tiếp theo là biến năng lực hiểu địa điểm này thành một agent đô thị thực tế: khách nói muốn đi đâu, agent lập lịch trình và hành động; người bán nói muốn mở mô hình gì, agent phân tích khu vực và giải thích cơ hội. Giá trị nằm ở lớp sản phẩm xây phía trên agent, không phải ở việc tự viết một agent rời rạc.
"""
    roadmap_path = OUT / "PISI_2026_dinh_huong_phat_trien_Danang_UrbanAgent_AI.md"
    roadmap_path.write_text(roadmap, encoding="utf-8")
    return roadmap_path


def main():
    OUT.mkdir(parents=True, exist_ok=True)
    stats = read_stats()
    outputs = [
        create_docx(stats),
        create_pptx(stats),
        *create_markdown_files(stats),
        create_roadmap_file(stats),
    ]
    print("\n".join(str(path) for path in outputs))


if __name__ == "__main__":
    main()
